from pptx import Presentation
import pdfkit
from PyPDF2 import PdfFileMerger, PdfFileReader
import os

# replace 'presentation.pptx' with your own .pptx filename
prs = Presentation('presenation.pptx')

output_folder = "temp"
if not os.path.exists(output_folder):
    os.makedirs(output_folder)

# write each slide to a separate pdf and store in the temp folder
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for text_frame in shape.text_frame.paragraphs:
               # write each line to the pdf
             text_run = text_frame.runs
             with open(f'{output_folder}/slide_{i}.pdf', 'wb') as f:
                 content = ''
                 for run in text_run:
                     content += run.text
                 f.write(content.encode('utf-8'))  # Write binary data instead of string
                 
# merge all slide PDFs into one single file
merger = PdfFileMerger()
for filename in os.listdir(output_folder):
    merger.append(PdfFileReader(f'{output_folder}/{filename}'))
    
with open("presentation.pdf", "wb") as outputStream:
    merger.write(outputStream)